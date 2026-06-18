from __future__ import annotations

import base64
import html
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageColor, ImageDraw, ImageFont


FONT_CANDIDATES = {
    "regular": (
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    ),
    "bold": (
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
    ),
}


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _attr_float(attrs: dict[str, str], key: str, default: float = 0.0) -> float:
    value = attrs.get(key)
    if value is None:
        return default
    value = value.strip().removesuffix("px")
    try:
        return float(value)
    except ValueError:
        return default


def _parse_svg_size(root: ET.Element) -> tuple[int, int]:
    view_box = root.attrib.get("viewBox") or root.attrib.get("viewbox")
    if view_box:
        parts = view_box.replace(",", " ").split()
        if len(parts) == 4:
            try:
                width = int(round(float(parts[2])))
                height = int(round(float(parts[3])))
                if width > 0 and height > 0:
                    return width, height
            except ValueError:
                pass
    width = int(round(_attr_float(root.attrib, "width", 1080)))
    height = int(round(_attr_float(root.attrib, "height", 1080)))
    return max(width, 1), max(height, 1)


def _rgba(fill: str | None, opacity: float = 1.0) -> tuple[int, int, int, int] | None:
    if not fill or fill == "none":
        return None
    try:
        red, green, blue = ImageColor.getrgb(fill)[:3]
    except ValueError:
        return None
    alpha = max(0, min(255, int(round(255 * opacity))))
    return red, green, blue, alpha


def _font(size: int, weight: str | None = None) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    family = "bold" if weight and (weight == "bold" or weight.isdigit() and int(weight) >= 700) else "regular"
    for path in FONT_CANDIDATES[family]:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, max(1, size))
            except OSError:
                continue
    return ImageFont.load_default()


def _alpha_draw(base: Image.Image) -> ImageDraw.ImageDraw:
    return ImageDraw.Draw(base, "RGBA")


def _draw_rect(base: Image.Image, attrs: dict[str, str]) -> None:
    fill = _rgba(attrs.get("fill"), _attr_float(attrs, "opacity", 1.0))
    if fill is None:
        return
    x = _attr_float(attrs, "x")
    y = _attr_float(attrs, "y")
    width = _attr_float(attrs, "width")
    height = _attr_float(attrs, "height")
    radius = _attr_float(attrs, "rx")
    draw = _alpha_draw(base)
    box = (x, y, x + width, y + height)
    if radius > 0:
        draw.rounded_rectangle(box, radius=radius, fill=fill)
    else:
        draw.rectangle(box, fill=fill)


def _draw_circle(base: Image.Image, attrs: dict[str, str]) -> None:
    fill = _rgba(attrs.get("fill"), _attr_float(attrs, "opacity", 1.0))
    if fill is None:
        return
    cx = _attr_float(attrs, "cx")
    cy = _attr_float(attrs, "cy")
    radius = _attr_float(attrs, "r")
    _alpha_draw(base).ellipse((cx - radius, cy - radius, cx + radius, cy + radius), fill=fill)


def _draw_text(base: Image.Image, node: ET.Element) -> None:
    text = "".join(node.itertext()).strip()
    if not text:
        return
    attrs = node.attrib
    fill = _rgba(attrs.get("fill"), _attr_float(attrs, "opacity", 1.0))
    if fill is None:
        return
    x = _attr_float(attrs, "x")
    y = _attr_float(attrs, "y")
    size = int(round(_attr_float(attrs, "font-size", 24)))
    font = _font(size, attrs.get("font-weight"))
    anchor = "ls"
    if attrs.get("text-anchor") == "middle":
        anchor = "ms"
    try:
        _alpha_draw(base).text((x, y), html.unescape(text), font=font, fill=fill, anchor=anchor)
    except TypeError:
        _alpha_draw(base).text((x, y - size), html.unescape(text), font=font, fill=fill)


def _href(attrs: dict[str, str]) -> str | None:
    for key, value in attrs.items():
        if key == "href" or key.endswith("}href"):
            return value
    return None


def _draw_image(base: Image.Image, attrs: dict[str, str]) -> None:
    href = _href(attrs)
    if not href or "base64," not in href:
        return
    try:
        raw = base64.b64decode(href.split("base64,", 1)[1])
        image = Image.open(BytesIO(raw)).convert("RGBA")
    except Exception:
        return
    x = _attr_float(attrs, "x")
    y = _attr_float(attrs, "y")
    width = _attr_float(attrs, "width")
    height = _attr_float(attrs, "height")
    if width <= 0 or height <= 0 or image.width <= 0 or image.height <= 0:
        return
    scale = min(width / image.width, height / image.height)
    resized = image.resize((max(1, int(image.width * scale)), max(1, int(image.height * scale))), Image.Resampling.LANCZOS)
    paste_x = int(round(x + (width - resized.width) / 2))
    paste_y = int(round(y + (height - resized.height) / 2))
    base.alpha_composite(resized, (paste_x, paste_y))


def poster_svg_to_png(svg: str) -> bytes:
    root = ET.fromstring(svg)
    width, height = _parse_svg_size(root)
    image = Image.new("RGBA", (width, height), (255, 255, 255, 0))

    def walk(node: ET.Element) -> None:
        for child in list(node):
            name = _local_name(child.tag)
            if name in {"defs", "clipPath"}:
                continue
            if name == "rect":
                _draw_rect(image, child.attrib)
            elif name == "circle":
                _draw_circle(image, child.attrib)
            elif name == "image":
                _draw_image(image, child.attrib)
            elif name == "text":
                _draw_text(image, child)
            walk(child)

    walk(root)
    output = BytesIO()
    image.convert("RGB").save(output, format="PNG", optimize=True)
    return output.getvalue()


def _hex_to_rgb(color: str) -> tuple[int, int, int]:
    red, green, blue = ImageColor.getrgb(color)[:3]
    return red, green, blue


def _fit_picture(slide, image_bytes: bytes, left: int, top: int, width: int, height: int) -> None:
    source = BytesIO(image_bytes)
    with Image.open(BytesIO(image_bytes)) as image:
        image_ratio = image.width / image.height
    target_ratio = width / height
    if image_ratio >= target_ratio:
        pic_width = width
        pic_height = int(width / image_ratio)
    else:
        pic_height = height
        pic_width = int(height * image_ratio)
    slide.shapes.add_picture(source, left + int((width - pic_width) / 2), top + int((height - pic_height) / 2), pic_width, pic_height)


def _apply_text_run(run, *, font_size, bold: bool = False, color: str = "#1f2937") -> None:
    from pptx.dml.color import RGBColor

    run.font.name = "Malgun Gothic"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*_hex_to_rgb(color))


def _add_text(slide, left, top, width, height, text: str, *, font_size, bold: bool = False, color: str = "#1f2937") -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    _apply_text_run(run, font_size=font_size, bold=bold, color=color)


def build_poster_pptx(*, poster_svg: str | None, copy_result: dict, poster: dict | None, product: dict) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(*_hex_to_rgb("#f6f8fc"))

    left = Inches(0.45)
    top = Inches(0.45)
    poster_w = Inches(7.25)
    poster_h = Inches(6.6)
    right = Inches(8.05)
    right_w = Inches(4.75)

    if poster_svg:
        try:
            _fit_picture(slide, poster_svg_to_png(poster_svg), left, top, poster_w, poster_h)
        except Exception:
            _add_text(slide, left, top, poster_w, Inches(0.45), "Poster SVG", font_size=Pt(18), bold=True)
            _add_text(slide, left, top + Inches(0.6), poster_w, Inches(5.7), poster_svg[:1800], font_size=Pt(8), color="#475569")

    headline = str(copy_result.get("headline") or product.get("product_name") or "")
    subcopy = str(copy_result.get("subcopy") or product.get("selling_point") or "")
    cta = str(copy_result.get("cta") or "자세히 보기")
    copies = [str(item) for item in (copy_result.get("copies") or [])[:4]]
    hashtags = " ".join(str(tag) for tag in (copy_result.get("hashtags") or [])[:5])

    _add_text(slide, right, Inches(0.62), right_w, Inches(0.32), str(product.get("target_channel") or "광고 포스터"), font_size=Pt(11), bold=True, color="#2563eb")
    _add_text(slide, right, Inches(1.05), right_w, Inches(0.95), headline, font_size=Pt(26), bold=True, color="#111827")
    _add_text(slide, right, Inches(2.1), right_w, Inches(0.85), subcopy, font_size=Pt(15), color="#475569")

    cta_units = sum(1.0 if ord(char) > 127 else 0.4 if char.isspace() else 0.6 for char in cta)
    cta_width = min(right_w, max(Inches(1.9), Inches(cta_units * 12 / 72 + 0.7)))
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right, Inches(3.1), cta_width, Inches(0.48))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = RGBColor(*_hex_to_rgb("#3b82f6"))
    cta_box.line.color.rgb = RGBColor(*_hex_to_rgb("#3b82f6"))
    cta_box.text_frame.clear()
    cta_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cta_paragraph = cta_box.text_frame.paragraphs[0]
    cta_paragraph.alignment = PP_ALIGN.CENTER
    cta_run = cta_paragraph.add_run()
    cta_run.text = cta
    _apply_text_run(cta_run, font_size=Pt(12), bold=True, color="#ffffff")

    y = Inches(3.9)
    for line in copies:
        _add_text(slide, right, y, right_w, Inches(0.34), f"- {line}", font_size=Pt(12), color="#334155")
        y += Inches(0.45)

    meta = f"{product.get('product_name', '')} · {product.get('price', '')}".strip(" ·")
    if meta:
        _add_text(slide, right, Inches(6.2), right_w, Inches(0.3), meta, font_size=Pt(10), color="#64748b")
    if hashtags:
        _add_text(slide, right, Inches(6.55), right_w, Inches(0.3), hashtags, font_size=Pt(9), color="#2563eb")

    if poster:
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"poster_template: {poster.get('poster_template', '')}\nposter_url: {poster.get('poster_url', '')}"

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
