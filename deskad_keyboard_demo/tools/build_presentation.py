#!/usr/bin/env python
"""DeskAd 발표 산출물 빌더 — 프로젝트 보고서(PDF) + 슬라이드(PPTX).

자산:
  docs/presentation/assets/img/*.png   예시 이미지(UI·3D·depth·광고)
  docs/presentation/assets/fonts/NanumGothic-*.ttf  한글 폰트(없으면 download_font로 받음)
출력:
  docs/presentation/DeskAd_프로젝트_보고서.pdf
  docs/presentation/DeskAd_발표.pptx

사용: conda run -n sprint_high python tools/build_presentation.py
모든 docs/* 는 기본 .gitignore로 로컬 전용(공개 정책). 인프라 값은 본문에서 마스킹.
"""
from __future__ import annotations

import urllib.request
from pathlib import Path

from PIL import Image as PILImage

ROOT = Path(__file__).resolve().parent.parent
PRES = ROOT / "docs" / "presentation"
IMG = PRES / "assets" / "img"
FONTS = PRES / "assets" / "fonts"
PDF_OUT = PRES / "DeskAd_프로젝트_보고서.pdf"
PPTX_OUT = PRES / "DeskAd_발표.pptx"

FONT_URLS = {
    "NanumGothic-Regular.ttf": "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Regular.ttf",
    "NanumGothic-Bold.ttf": "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Bold.ttf",
}


def ensure_fonts() -> tuple[Path, Path]:
    FONTS.mkdir(parents=True, exist_ok=True)
    for name, url in FONT_URLS.items():
        p = FONTS / name
        if not p.exists():
            print("downloading font", name)
            urllib.request.urlretrieve(url, p)
    return FONTS / "NanumGothic-Regular.ttf", FONTS / "NanumGothic-Bold.ttf"


def img_size(name: str) -> tuple[int, int] | None:
    p = IMG / name
    if not p.exists():
        return None
    with PILImage.open(p) as im:
        return im.size


# ── 색상 팔레트 ───────────────────────────────────────────────────────────────
INK = "#1f2430"
ACCENT = "#c8552e"
MUTED = "#5a616b"
PANEL = "#f4f1eb"
BLUE = "#3766c4"


# ============================================================================
# 1) PDF 보고서 (reportlab Platypus)
# ============================================================================
def build_pdf() -> None:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    reg, bold = ensure_fonts()
    pdfmetrics.registerFont(TTFont("Nanum", str(reg)))
    pdfmetrics.registerFont(TTFont("Nanum-Bold", str(bold)))
    pdfmetrics.registerFontFamily("Nanum", normal="Nanum", bold="Nanum-Bold")

    from reportlab.platypus import (
        Image,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    styles = getSampleStyleSheet()

    def st(name, **kw):
        base = dict(fontName="Nanum", textColor=colors.HexColor(INK), leading=16)
        base.update(kw)
        return ParagraphStyle(name, parent=styles["Normal"], **base)

    body = st("body", fontSize=10.2, leading=16, spaceAfter=6)
    h1 = st("h1", fontName="Nanum-Bold", fontSize=18, textColor=colors.HexColor(ACCENT),
            spaceBefore=14, spaceAfter=8, leading=22)
    h2 = st("h2", fontName="Nanum-Bold", fontSize=13, textColor=colors.HexColor(INK),
            spaceBefore=10, spaceAfter=5, leading=17)
    cap = st("cap", fontSize=8.4, textColor=colors.HexColor(MUTED), alignment=TA_CENTER,
             spaceBefore=3, spaceAfter=10, leading=11)
    bullet = st("bullet", fontSize=10.2, leading=15, leftIndent=12, spaceAfter=3, bulletIndent=2)
    title = st("title", fontName="Nanum-Bold", fontSize=30, alignment=TA_LEFT,
               textColor=colors.HexColor(INK), leading=36)
    subtitle = st("subtitle", fontSize=13, textColor=colors.HexColor(MUTED), leading=20)

    story: list = []

    def para(text, style=body):
        story.append(Paragraph(text, style))

    def bullets(items):
        for it in items:
            story.append(Paragraph(f"•&nbsp;&nbsp;{it}", bullet))

    def figure(name, caption, max_w=15.5):
        sz = img_size(name)
        if not sz:
            return
        w, h = sz
        draw_w = min(max_w * cm, 16 * cm)
        draw_h = draw_w * h / w
        max_h = 11 * cm
        if draw_h > max_h:
            draw_h = max_h
            draw_w = draw_h * w / h
        story.append(Image(str(IMG / name), width=draw_w, height=draw_h))
        story.append(Paragraph(caption, cap))

    def figure_row(pairs, max_w=15.8):
        # pairs: [(name, caption), ...] 가로 배치
        cells, capcells = [], []
        cw = (max_w * cm) / len(pairs) - 0.25 * cm
        for name, caption in pairs:
            sz = img_size(name)
            if not sz:
                cells.append("")
                capcells.append("")
                continue
            w, h = sz
            dw = cw
            dh = dw * h / w
            cells.append(Image(str(IMG / name), width=dw, height=dh))
            capcells.append(Paragraph(caption, cap))
        t = Table([cells, capcells], colWidths=[cw + 0.25 * cm] * len(pairs))
        t.setStyle(TableStyle([("ALIGN", (0, 0), (-1, -1), "CENTER"),
                               ("VALIGN", (0, 0), (-1, -1), "TOP")]))
        story.append(t)
        story.append(Spacer(1, 0.3 * cm))

    def table(rows, header=True, col_widths=None):
        data = [[Paragraph(str(c), st("cell", fontSize=9, leading=12,
                                      fontName="Nanum-Bold" if (header and ri == 0) else "Nanum",
                                      textColor=colors.white if (header and ri == 0) else colors.HexColor(INK)))
                 for c in row] for ri, row in enumerate(rows)]
        t = Table(data, colWidths=col_widths, hAlign="LEFT")
        ts = [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(INK)),
              ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#d8d3c8")),
              ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f7f5f0")]),
              ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
              ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
              ("LEFTPADDING", (0, 0), (-1, -1), 6)]
        t.setStyle(TableStyle(ts))
        story.append(t)
        story.append(Spacer(1, 0.35 * cm))

    # ── 표지 ──
    story.append(Spacer(1, 3.2 * cm))
    para("DeskAd AI Studio", title)
    story.append(Spacer(1, 0.3 * cm))
    para("3D 데스크 셋업 미리보기 + 배열 충실도 기반 광고 콘텐츠 자동 생성", subtitle)
    story.append(Spacer(1, 0.5 * cm))
    para("프로젝트 보고서", st("ptag", fontName="Nanum-Bold", fontSize=14,
                          textColor=colors.HexColor(ACCENT)))
    story.append(Spacer(1, 0.2 * cm))
    para("작성일 2026-06-17 · 소상공인 커스텀 키보드/데스크테리어 판매자용", subtitle)
    story.append(Spacer(1, 0.8 * cm))
    figure("setup_3d_render.png", "절차적 3D 셋업 렌더(65% 키보드·데스크매트·모니터·마우스) — 광고 생성의 구조 레퍼런스", max_w=13)
    story.append(PageBreak())

    # ── 1. 개요 ──
    para("1. 프로젝트 개요", h1)
    para("DeskAd AI Studio는 커스텀 키보드와 데스크테리어를 파는 소상공인이 <b>상품 스펙과 데스크 환경만 입력하면 "
         "3D 셋업 미리보기와 사진 품질의 광고(문구+이미지+포스터)를 자동으로 얻도록</b> 돕는 도구다. "
         "디자이너·촬영 스튜디오 없이도 채널(인스타그램·스마트스토어·상세페이지 등)에 맞는 광고 소재를 만들 수 있다.")
    para("핵심 차별점은 <b>“AI가 그린 이미지가 실제 판매 제품과 다르다”는 생성형 광고의 근본 문제를, "
         "생성한 3D 셋업을 구조 레퍼런스로 삼는 배열 충실도 고정(depth-ControlNet)으로 해결</b>한 점이다(5장).")

    para("2. 배경과 문제 정의", h1)
    para("소상공인은 광고 이미지 제작에 디자인·촬영 비용과 시간을 들이기 어렵다. 범용 이미지 생성 모델을 쓰면 빠르지만, "
         "키보드 같은 정밀 제품에서는 <b>배열이 틀리거나(65%를 풀사이즈로), 행이 물결치거나, 키캡이 뭉개지는</b> 현상이 "
         "분산적으로 발생한다. 즉 ‘그럴듯하지만 실제 상품과 다른’ 이미지가 나와 광고로 쓸 수 없다.")
    para("DeskAd는 (1) 판매자의 실제 스펙으로 정확한 3D 셋업을 만들고, (2) 그 구조를 생성 단계에 강제로 주입해 "
         "<b>사진 품질과 정확한 배열을 동시에</b> 얻는 방식으로 이 문제를 정면 해결한다.")

    # ── 3. 시스템 개요 ──
    para("3. 시스템 개요", h1)
    para("브라우저 → (nginx + Basic Auth) → Streamlit UI(내부) → FastAPI 백엔드(내부) → 3D 렌더 / 텍스트 LLM / 이미지 백엔드. "
         "백엔드가 모든 입력 검증·프롬프트 구성·잡 큐·캐시·GPU 워커 수명주기를 담당하는 단일 오케스트레이터다.")
    para("이미지/텍스트 엔진은 2트랙이다.", h2)
    table([["엔진", "광고 문구", "광고 이미지", "특징"],
           ["local", "로컬 LLM(HyperCLOVA X SEED·Kanana·Mi:dm)", "ComfyUI/FLUX (depth-ControlNet/img2img)", "키 불필요, 배열 정확도 우선"],
           ["openai", "OpenAI 호환 API", "OpenAI 이미지", "OPENAI_API_KEY 필요"]],
          col_widths=[1.8 * cm, 5.4 * cm, 5.2 * cm, 4.0 * cm])
    para("GPU 워커(ComfyUI/HyperCLOVA)는 단일 NVIDIA L4(24GB)를 공유하며 <b>GPU_WORKER_MODE</b>(always_on·exclusive·on_demand)로 "
         "수명주기를 제어한다. 엔진/워커/키가 없으면 문구는 템플릿, 이미지는 SVG 일러스트로 안전하게 폴백한다.")
    figure("ui_3d_preview.png", "실제 UI — ③ 3D 미리보기 단계(스펙 입력 → GLB 생성)", max_w=14)

    # ── 4. 사용 흐름 ──
    story.append(PageBreak())
    para("4. 사용 흐름 — 4단계 위저드", h1)
    table([["단계", "입력", "처리", "산출"],
           ["① 상품 정보", "모델명·레이아웃, STEP/STP/GLB 업로드", "업로드 검증·변환", "정규화 상품 메타"],
           ["② 도면 미리보기", "레이아웃 JSON", "→ SVG 탑뷰 도면", "배열 확인 도면"],
           ["③ 3D 셋업", "색·재질·책상/모니터/소품", "절차적 GLB 빌드(1 unit=1cm)", "GLB + 구도 맵"],
           ["④ 광고 생성", "채널·톤·셀링포인트", "문구(LLM)+이미지(엔진)+포스터", "카피·이미지·SVG/PPTX"]],
          col_widths=[2.6 * cm, 4.4 * cm, 4.8 * cm, 4.0 * cm])
    para("3D 셋업 빌더는 GLB 단위=1cm 규약과 MX 표준 간격(19.05mm)을 지켜 키보드 footprint·키캡 프로파일·책상 치수를 "
         "실측 기반으로 생성한다. 이 정확한 3D가 다음 장의 충실도 고정의 토대가 된다. 빌더는 동시에 광고용 ‘구도 맵’(평면 색블록)을 "
         "원근/탑다운 두 투영으로 함께 만든다.")
    figure_row([("composition_perspective.png", "구도 맵(원근)"),
                ("composition_topdown.png", "구도 맵(탑다운/flat-lay)")])

    # ── 5. 핵심 기술 ──
    story.append(PageBreak())
    para("5. 핵심 기술 — 배열 충실도 (depth-ControlNet)", h1)
    para("순수 text2image나 평면 도면 img2img로는 ‘사진 품질’과 ‘정확한 배열’을 동시에 얻기 어렵다. "
         "DeskAd는 <b>구조와 외관을 분리</b>해 이 문제를 푼다.")
    para("셋업 GLB → (OSMesa, CPU 헤드리스 렌더) → depth PNG → FLUX depth-ControlNet → 사진 광고", h2)
    para("생성한 3D 셋업 GLB를 소프트웨어 렌더(OSMesa, <b>GPU 미사용</b>)로 grayscale depth로 만들고, ComfyUI의 depth-ControlNet이 "
         "이 depth로 <b>구조를 denoise와 독립적으로 고정</b>한다. 그래서 denoise=1.0(완전한 사진 자유도)이어도 65% 배열이 무너지지 않는다.")
    figure_row([("composition_perspective.png", "① 셋업(구도 맵)"),
                ("depth_hero.png", "② depth 입력"),
                ("ad_hero.png", "③ 최종 사진 광고")])
    para("색은 depth가 grayscale라 고정되지 않으므로 프롬프트 그라운딩과 best-of-N이 분담한다. "
         "즉 <b>배열=depth · 주색=프롬프트 · 액센트=best-of-N</b>의 3단 직교 구성이다.", body)
    table([["노브(.env)", "역할"],
           ["COMFYUI_CONTROLNET_STRENGTH", "구조 강제 강도(0.5=스위트스팟, 0.7↑ 평면화, 0=비활성→img2img)"],
           ["COMFYUI_CONTROLNET_END_PERCENT", "초기 스텝만 적용(<1.0이면 후반 사진 자유도↑)"],
           ["COMFYUI_BEST_OF_N", "N장 중 액센트 색 충실도 최적 컷 선택(단일 L4 권장 2~4)"]],
          col_widths=[6.0 * cm, 9.6 * cm])
    para("depth는 세워진 모니터를 포함한 데스크 시점이라 <b>데스크/룸 구도</b>(hero·eye_level·wide)에만 쓰고, flat-lay·매크로 컷은 "
         "ControlNet을 끄고 img2img로 자연 폴백한다.")

    # ── 6. 컷별 시점 + 순차 ──
    story.append(PageBreak())
    para("6. 그리드 광고 — 컷별 시점 분리 + 순차 생성", h1)
    para("grid_three 포스터는 hero(메인)·detail_macro(키캡 디테일)·eye_level(데스크 무드) <b>3컷을 서로 다른 시점</b>으로 만든다. "
         "depth 컷(hero·eye_level)은 같은 GLB를 <b>컷별 카메라 각도</b>(구면좌표 프리셋: hero=높은 3/4, eye_level=낮은 수평)로 렌더해 "
         "시점이 실제로 갈린다.")
    figure_row([("depth_hero.png", "depth: hero(높은 3/4)"),
                ("depth_eye_level.png", "depth: eye_level(낮은 수평)"),
                ("depth_legacy.png", "(이전) 고정 각도 — 두 컷 공유")])
    para("이전에는 두 컷이 위 오른쪽의 ‘고정 각도’ depth 하나를 공유해 시점이 겹쳤다. 카메라를 컷별로 분리하면서, 컷별 depth를 ComfyUI에 "
         "<b>유니크 파일명</b>으로 올리도록 함께 고쳤다(고정 파일명+overwrite는 LoadImage가 실행 시점에 파일을 읽는 탓에 서로 덮어써 "
         "시점이 다시 겹치는 함정). 그 결과 최종 FLUX 이미지에서도 시점이 명확히 갈린다.")
    figure_row([("ad_hero.png", "최종: hero"),
                ("ad_eye_level.png", "최종: eye_level"),
                ("ad_detail_macro.png", "최종: detail_macro")])
    para("순차 생성 (안정성)", h2)
    para("그리드 3컷을 ComfyUI에 한꺼번에 큐잉하면 단일 L4에서 FLUX+ControlNet 컷의 VRAM 피크가 겹쳐 서버가 죽을 수 있다. "
         "그래서 <b>첫 컷만 제출하고, 폴링이 현재 컷 완료를 확인한 뒤에야 다음 컷을 제출</b>하도록 바꿔 ComfyUI 큐에 우리 컷이 "
         "항상 1개만 존재하게 했다. 라이브 검증에서 매 순간 in-flight=1로 크래시 없이 3컷을 완주했다.")

    # ── 7. 보안 ──
    story.append(PageBreak())
    para("7. 보안 요약", h1)
    para("상세는 docs/security.md. 위협 모델 우선순위에 따라 다층 방어를 적용했다.")
    bullets([
        "<b>시크릿 위생</b>: API 키·토큰을 화면/로그/git에 노출하지 않음(마스킹), pre-commit 시크릿 스캔.",
        "<b>접근 제어</b>: Streamlit·모델 워커 포트를 외부에 직접 열지 않고 nginx + Basic Auth + 세션 인증 뒤에 둠.",
        "<b>프롬프트 인젝션</b>: 한/영 인젝션 패턴 탐지(감사 플래그) + 시스템 프롬프트 가드레일.",
        "<b>입력 검증</b>: 모든 요청이 Pydantic 스키마(길이·패턴·범위)를 통과, 업로드 GLB/경로 탈출 차단(basename만 신뢰).",
        "<b>리소스</b>: GPU 워커 수명주기·요청 타임아웃·결과 캐시로 고갈 완화, 런타임 상태/락 파일 0600.",
    ])

    # ── 8. 품질 ──
    para("8. 품질 보증과 성과", h1)
    bullets([
        "회귀 테스트 <b>263개</b> 통과(conda run -n sprint_high pytest).",
        "충실도·구도 변경은 단위 테스트로 동작을 고정하고, 시각 품질(각도·사진감)은 <b>라이브 ComfyUI 실생성으로 눈 검증</b>.",
        "배열 충실도: depth-ControlNet strength 0.5에서 ‘사진 + 65% 정확 배열’ 동시 달성(라이브).",
        "그리드 순차 생성: in-flight=1 유지, 3컷 크래시 없이 완주(라이브).",
        "백엔드 변경은 라이브 검증 전 start.sh --restart로 코드 반영(uvicorn 기동 시점 고정).",
    ])

    para("9. 향후 과제", h1)
    bullets([
        "레이아웃 인지 VLM judge(GPT-4V급)로 best-of-N 의미 점수 강화 — 넘패드 오배열·행 붕괴까지 자동 강등.",
        "포스터 레이아웃 추가 다듬기(긴 제품명/가격 오버플로, 테마별 대비).",
        "문구→이미지→포스터를 하나의 파이프라인 단계 게이지로 통합.",
    ])

    doc = SimpleDocTemplate(str(PDF_OUT), pagesize=A4,
                            leftMargin=2.2 * cm, rightMargin=2.2 * cm,
                            topMargin=2.0 * cm, bottomMargin=1.8 * cm,
                            title="DeskAd AI Studio 프로젝트 보고서", author="DeskAd")

    def footer(canvas, d):
        canvas.saveState()
        canvas.setFont("Nanum", 8)
        canvas.setFillColor(colors.HexColor(MUTED))
        canvas.drawString(2.2 * cm, 1.1 * cm, "DeskAd AI Studio · 프로젝트 보고서")
        canvas.drawRightString(A4[0] - 2.2 * cm, 1.1 * cm, "%d" % d.page)
        canvas.restoreState()

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    print("PDF ->", PDF_OUT)


# ============================================================================
# 2) PPTX 슬라이드 (python-pptx)
# ============================================================================
def build_pptx() -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    def rgb(h):
        return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

    KOR = "맑은 고딕"  # PowerPoint 한글 렌더용(없으면 시스템 한글 폰트로 대체)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    def rect(s, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(color)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(s, x, y, w, h, lines, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=KOR):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        if isinstance(lines, str):
            lines = [lines]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if isinstance(ln, tuple):
                text, lvl, sz, bd = ln
            else:
                text, lvl, sz, bd = ln, 0, size, bold
            p.level = lvl
            run = p.add_run()
            run.text = text
            run.font.size = Pt(sz)
            run.font.bold = bd
            run.font.name = font
            run.font.color.rgb = rgb(color if not isinstance(ln, tuple) else INK)
        return tb

    def bullets_box(s, x, y, w, h, items, size=16):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            lvl = 0
            txt = it
            if isinstance(it, tuple):
                txt, lvl = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = lvl
            run = p.add_run()
            run.text = ("• " if lvl == 0 else "– ") + txt
            run.font.size = Pt(size - lvl * 2)
            run.font.name = KOR
            run.font.color.rgb = rgb(INK if lvl == 0 else MUTED)
            p.space_after = Pt(6)
        return tb

    def pic_fit(s, name, x, y, max_w, max_h):
        sz = img_size(name)
        if not sz:
            return
        w, h = sz
        scale = min(max_w / w, max_h / h)
        dw, dh = int(w * scale), int(h * scale)
        s.shapes.add_picture(str(IMG / name), x + (max_w - dw) // 2, y + (max_h - dh) // 2, dw, dh)

    def header(s, title, idx=None):
        rect(s, 0, 0, SW, Inches(1.15), INK)
        rect(s, 0, Inches(1.15), SW, Pt(4), ACCENT)
        textbox(s, Inches(0.55), Inches(0.18), SW - Inches(1.1), Inches(0.85),
                title, size=26, color="#ffffff", bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if idx:
            textbox(s, SW - Inches(1.4), Inches(0.18), Inches(0.9), Inches(0.85),
                    idx, size=13, color="#c8c1b2", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ── 1. 표지 ──
    s = slide()
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, Inches(4.05), SW, Pt(5), ACCENT)
    textbox(s, Inches(0.8), Inches(2.2), SW - Inches(1.6), Inches(1.4),
            "DeskAd AI Studio", size=52, color="#ffffff", bold=True)
    textbox(s, Inches(0.85), Inches(4.25), SW - Inches(1.6), Inches(1.6),
            ["3D 데스크 셋업 미리보기 + 배열 충실도 기반 광고 콘텐츠 자동 생성",
             "소상공인 커스텀 키보드 · 데스크테리어 판매자용 · 2026-06-17"],
            size=18, color="#c8c1b2")

    # ── 2. 목차 ──
    s = slide()
    header(s, "목차", "01")
    bullets_box(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(5.6), [
        "배경과 문제 — ‘AI 광고가 실제 제품과 다르다’",
        "시스템 개요 — UI · 백엔드 · 2트랙 엔진",
        "사용 흐름 — 4단계 위저드",
        "핵심 기술 — 배열 충실도(depth-ControlNet)",
        "그리드 광고 — 컷별 시점 분리 + 순차 생성",
        "보안 · 품질 · 성과",
    ], size=20)

    # ── 3. 문제 ──
    s = slide()
    header(s, "배경과 문제 정의", "02")
    bullets_box(s, Inches(0.7), Inches(1.5), Inches(6.4), Inches(5.4), [
        "소상공인은 광고 촬영·디자인에 비용·시간을 들이기 어렵다.",
        "범용 생성 모델은 빠르지만 정밀 제품에서 배열이 틀린다:",
        ("65% 키보드를 풀사이즈로 생성", 1),
        ("행 물결침 · 키캡 뭉개짐(melt)", 1),
        "→ ‘그럴듯하지만 실제 상품과 다른’ 이미지 = 광고로 못 씀.",
        "DeskAd: 실제 스펙의 3D를 구조로 강제 주입해 해결.",
    ], size=17)
    pic_fit(s, "ad_hero.png", Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.0))

    # ── 4. 시스템 개요 ──
    s = slide()
    header(s, "시스템 개요", "03")
    bullets_box(s, Inches(0.7), Inches(1.45), Inches(6.5), Inches(5.6), [
        "브라우저 → nginx(Basic Auth) → Streamlit UI → FastAPI 백엔드",
        "백엔드 = 단일 오케스트레이터(검증·프롬프트·잡·캐시·워커)",
        "이미지/텍스트 2트랙:",
        ("local: 로컬 LLM + ComfyUI/FLUX (정확도 우선, 키 불필요)", 1),
        ("openai: OpenAI 텍스트 + 이미지 (키 필요)", 1),
        "GPU 워커는 단일 L4(24GB) 공유 — GPU_WORKER_MODE로 제어",
        "엔진/키 없으면 템플릿·SVG로 안전 폴백",
    ], size=16)
    pic_fit(s, "ui_3d_preview.png", Inches(7.45), Inches(1.7), Inches(5.3), Inches(4.8))

    # ── 5. 4단계 ──
    s = slide()
    header(s, "사용 흐름 — 4단계 위저드", "04")
    bullets_box(s, Inches(0.7), Inches(1.45), Inches(6.4), Inches(5.6), [
        "① 상품 정보 — 모델명·레이아웃, STEP/GLB 업로드",
        "② 도면 미리보기 — 레이아웃 → SVG 탑뷰",
        "③ 3D 셋업 — 색·재질·책상/모니터/소품 → GLB(1 unit=1cm)",
        "④ 광고 생성 — 문구 + 이미지 + SVG/PPTX 포스터",
        ("MX 표준 간격 19.05mm·실측 footprint → 충실도의 토대", 1),
    ], size=17)
    pic_fit(s, "setup_3d_render.png", Inches(7.4), Inches(1.6), Inches(5.4), Inches(5.2))

    # ── 6. 핵심 기술 ──
    s = slide()
    header(s, "핵심 기술 — 배열 충실도 (depth-ControlNet)", "05")
    textbox(s, Inches(0.7), Inches(1.3), SW - Inches(1.4), Inches(0.9),
            "셋업 GLB →(CPU 헤드리스 렌더) depth →  FLUX depth-ControlNet → 사진 광고  ·  구조와 외관을 분리",
            size=16, color=ACCENT, bold=True)
    y = Inches(2.3)
    for i, (name, capt) in enumerate([("composition_perspective.png", "① 셋업(구도 맵)"),
                                      ("depth_hero.png", "② depth 입력"),
                                      ("ad_hero.png", "③ 최종 사진 광고")]):
        x = Inches(0.7 + i * 4.25)
        pic_fit(s, name, x, y, Inches(3.9), Inches(3.6))
        textbox(s, x, Inches(5.95), Inches(3.9), Inches(0.5), capt, size=14,
                color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    textbox(s, Inches(0.7), Inches(6.55), SW - Inches(1.4), Inches(0.7),
            "배열=depth · 주색=프롬프트 · 액센트=best-of-N (3단 직교) · strength 0.5 = 사진+정확 배열 스위트스팟",
            size=14, color=INK)

    # ── 7. 컷별 시점 ──
    s = slide()
    header(s, "그리드 광고 — 컷별 시점 분리", "06")
    for i, (name, capt) in enumerate([("depth_hero.png", "depth: hero(높은 3/4)"),
                                      ("depth_eye_level.png", "depth: eye_level(낮은 수평)"),
                                      ("depth_legacy.png", "(이전) 고정 각도 — 겹침")]):
        x = Inches(0.7 + i * 4.25)
        pic_fit(s, name, x, Inches(1.4), Inches(3.9), Inches(2.5))
        textbox(s, x, Inches(3.85), Inches(3.9), Inches(0.4), capt, size=12,
                color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    for i, (name, capt) in enumerate([("ad_hero.png", "최종 hero"),
                                      ("ad_eye_level.png", "최종 eye_level"),
                                      ("ad_detail_macro.png", "최종 detail_macro")]):
        x = Inches(0.7 + i * 4.25)
        pic_fit(s, name, x, Inches(4.35), Inches(3.9), Inches(2.4))
        textbox(s, x, Inches(6.75), Inches(3.9), Inches(0.4), capt, size=12,
                color=MUTED, align=PP_ALIGN.CENTER, bold=True)

    # ── 8. 순차 생성 ──
    s = slide()
    header(s, "그리드 — 순차 생성(안정성)", "07")
    bullets_box(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(5.5), [
        "문제: 3컷을 한꺼번에 큐잉 → 단일 L4 VRAM 피크 겹침 → 서버 다운 위험",
        "해법: 첫 컷만 제출 → 폴링이 완료 확인 후에야 다음 컷 제출",
        ("ComfyUI 큐에 우리 컷이 항상 1개만 존재", 1),
        ("컷별 depth는 유니크 파일명으로 업로드(overwrite 클로버 방지)", 1),
        "라이브 검증: in-flight=1 유지, 크래시 없이 3컷 완주",
        "회귀 테스트 263개 통과 + 라이브 눈 검증",
    ], size=18)

    # ── 9. 보안·품질 ──
    s = slide()
    header(s, "보안 · 품질 · 성과", "08")
    bullets_box(s, Inches(0.7), Inches(1.45), Inches(6.4), Inches(5.6), [
        "보안(docs/security.md):",
        ("시크릿 마스킹 + pre-commit 스캔", 1),
        ("nginx Basic Auth + 세션 인증, 워커 포트 비공개", 1),
        ("프롬프트 인젝션 탐지 + Pydantic 입력 검증", 1),
        "품질:",
        ("회귀 263개 통과 · 라이브 시각 검증", 1),
        ("strength 0.5 = 사진+정확 배열 동시 달성", 1),
    ], size=16)
    pic_fit(s, "ad_eye_level.png", Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.0))

    # ── 10. 마무리 ──
    s = slide()
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, Inches(3.7), SW, Pt(5), ACCENT)
    textbox(s, Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(1.2),
            "정확한 3D → 정확한 광고", size=40, color="#ffffff", bold=True)
    textbox(s, Inches(0.85), Inches(3.95), SW - Inches(1.6), Inches(1.2),
            ["DeskAd AI Studio — 배열 충실도로 ‘쓸 수 있는’ 생성형 광고",
             "감사합니다."], size=20, color="#c8c1b2")

    prs.save(str(PPTX_OUT))
    print("PPTX ->", PPTX_OUT, f"({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    PRES.mkdir(parents=True, exist_ok=True)
    build_pdf()
    build_pptx()
