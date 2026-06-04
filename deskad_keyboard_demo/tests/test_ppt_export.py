from io import BytesIO
from zipfile import ZipFile

from PIL import Image

from ppt_export import build_poster_pptx, poster_svg_to_png


def test_poster_svg_to_png_renders_basic_shapes_and_text():
    svg = (
        '<svg viewBox="0 0 320 180" xmlns="http://www.w3.org/2000/svg">'
        '<rect width="320" height="180" fill="#f5f2ea"/>'
        '<rect x="20" y="30" width="120" height="70" rx="12" fill="#8aa0a8"/>'
        '<text x="24" y="140" font-size="24" font-weight="800" fill="#2f3438">테스트 카피</text>'
        "</svg>"
    )

    png = poster_svg_to_png(svg)
    image = Image.open(BytesIO(png))

    assert image.size == (320, 180)
    assert image.getpixel((30, 40)) != image.getpixel((5, 5))


def test_build_poster_pptx_embeds_rendered_poster_image():
    svg = (
        '<svg viewBox="0 0 320 180" xmlns="http://www.w3.org/2000/svg">'
        '<rect width="320" height="180" fill="#ffffff"/>'
        '<rect x="24" y="24" width="120" height="80" fill="#3b82f6"/>'
        "</svg>"
    )
    pptx = build_poster_pptx(
        poster_svg=svg,
        copy_result={
            "headline": "책상이 좋아지는 시간",
            "subcopy": "조용한 타건감과 정돈된 무드",
            "cta": "구성 보기",
            "copies": ["작은 책상에도 균형 있게 어울립니다"],
            "hashtags": ["#DeskSetup"],
        },
        poster={"poster_template": "minimal_card", "poster_url": "http://example.test/poster.svg"},
        product={"product_name": "Neo65", "price": "189,000원", "target_channel": "인스타그램"},
    )

    with ZipFile(BytesIO(pptx)) as archive:
        names = archive.namelist()

    assert "[Content_Types].xml" in names
    assert any(name.startswith("ppt/media/image") and name.endswith(".png") for name in names)


def test_cta_shape_width_scales_with_label_length():
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def cta_width(cta: str) -> int:
        pptx = build_poster_pptx(
            poster_svg=None,
            copy_result={"headline": "헤드라인", "subcopy": "서브카피", "cta": cta},
            poster=None,
            product={"product_name": "Neo65"},
        )
        shapes = Presentation(BytesIO(pptx)).slides[0].shapes
        return next(s.width for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)

    # A long CTA must produce a wider pill than a short one (no fixed-width clipping).
    assert cta_width("지금 30% 할인가로 만나보기 바로 지금 확인하세요") > cta_width("구매")
